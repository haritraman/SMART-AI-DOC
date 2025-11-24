from pptx import Presentation

def build_pptx(project, sections, filepath: str):
    prs = Presentation()

    for sec in sections:
        layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = sec.title

        content = sec.current_content or ""
        lines = [l.strip() for l in content.split("\n") if l.strip()]

        body = slide.placeholders[1]
        tf = body.text_frame

        first = True
        for line in lines:
            if first:
                tf.text = line
                first = False
            else:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0

    prs.save(filepath)
    return filepath
